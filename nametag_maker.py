import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

import copy
from pptx.dml.color import RGBColor

def convert_scheme_color_to_rgb(theme, scheme_color):
    # 테마에서 색상 팔레트를 가져옵니다.
    color_palette = theme.theme_part.theme_element.get_or_add_theme_colors().indexed_colors

    # 색상 인덱스를 가져옵니다.
    color_index = scheme_color.theme_color.value

    # 색상 팔레트에서 해당 색상을 찾습니다.
    color = color_palette[color_index]

    # RGB 값을 가져옵니다.
    rgb = (color.rgb.red, color.rgb.green, color.rgb.blue)

    # RGBColor 객체를 생성하여 반환합니다.
    return RGBColor(*rgb)


def read_pptx(ppt):
    slide_contents = []

    for slide in ppt.slides:
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)

        slide_contents.append(slide_text)

    return slide_contents

def copy_slide(prs):
    source_slide = prs.slides[0]
    slide_layout = prs.slide_layouts[6]
    copy_slide = prs.slides.add_slide(slide_layout)
    
    for shape in source_slide.shapes:
        el = shape.element
        newel = copy.deepcopy(el)
        copy_slide.shapes._spTree.insert_element_before(newel,'p:extLst')
    return prs

csv = pd.read_csv("/Users/kite/Desktop/2024 1팀/겨울수련회/1팀 겨울 리트릿.csv", encoding='euc-kr')
ppt = Presentation("/Users/kite/Desktop/2024 1팀/겨울수련회/ppt_sample.pptx")

for i in range(len(csv)-1):
    ppt = copy_slide(ppt)
    
slides = read_pptx(ppt)

for slide_index, slide_content in enumerate(slides):
    for content in slide_content:
        if content in csv.columns:
            slide_text = str(csv[content][slide_index])  # 값이 정수형일 경우 문자열로 변환
            for shape in ppt.slides[slide_index].shapes:
                if hasattr(shape, "text") and shape.text == content:
                    # 폰트 정보를 가져옴
                    font_style = shape.text_frame.paragraphs[0].runs[0].font  
                    shape.text = slide_text  # 기존 텍스트를 변경
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = font_style.bold
                            run.font.italic = font_style.italic
                            run.font.size = font_style.size
                            # 색깔은 일단 무조건 검은색 
                            run.font.name = font_style.name
            print(f"Changed content '{content}' in slide {slide_index + 1} to '{slide_text}'")

ppt.save('/Users/kite/Desktop/2024 1팀/겨울수련회/test.pptx')
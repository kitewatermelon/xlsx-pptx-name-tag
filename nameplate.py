import pandas as pd
from pptx import Presentation
import copy
import os
import re

def read_pptx(ppt):
    slide_contents = []

    for slide in ppt.slides:
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)

        slide_contents.append(slide_text)

    return slide_contents

def copy_slide_from_external_prs(prs):
    source_slide = prs.slides[0]
    slide_layout = prs.slide_layouts[6]
    curr_slide = prs.slides.add_slide(slide_layout)

    imgDict = {}

    for shp in source_slide.shapes:

        if re.search(r'그림 \d+', shp.name):
            # save image
            with open(shp.name+'.jpg', 'wb') as f:
                f.write(shp.image.blob)

            imgDict[shp.name+'.jpg'] = [shp.left, shp.top, shp.width, shp.height]
            
        else:
            el = shp.element
            newel = copy.deepcopy(el)

            curr_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for k, v in imgDict.items():
        # print(v[0], v[1], v[2], v[3])
        curr_slide.shapes.add_picture(k, v[0], v[1], v[2], v[3])
        os.remove(k)
        
    return prs


def make_nameplate(excel_url, pptx_url):
    try:
        csv = pd.read_excel(excel_url, engine='openpyxl')
    except:
        csv = pd.read_csv(excel_url, encoding='euc-kr')
    ppt = Presentation(pptx_url)

    for i in range(len(csv)-1):
        ppt = copy_slide_from_external_prs(ppt)
    slides = read_pptx(ppt)

    for slide_index, slide_content in enumerate(slides):
        for content in slide_content:
            if content in csv.columns:
                slide_text = str(csv[content][slide_index])
                for shape in ppt.slides[slide_index].shapes:
                    if hasattr(shape, "text") and shape.text == content:
                        font_style = shape.text_frame.paragraphs[0].runs[0].font  
                        shape.text = slide_text
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.bold = font_style.bold
                                run.font.italic = font_style.italic
                                run.font.size = font_style.size
                                run.font.name = font_style.name
                # print(f"Changed content '{content}' in slide {slide_index + 1} to '{slide_text}'")

    return ppt